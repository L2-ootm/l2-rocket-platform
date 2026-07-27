#!/usr/bin/env python3
"""
BatteryHive Investor Pitch Deck — 10 slides, data-heavy, premium design.
Style: Warm Modern (human, approachable, premium quality).
Generates .pptx with charts, tables, financial projections.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE
import os

# === DESIGN TOKENS (Warm Modern) ===
BG_DARK    = RGBColor(0x0F, 0x0E, 0x0C)  # near-black warm
BG_CARD    = RGBColor(0x1A, 0x18, 0x15)  # dark card
BG_SURFACE = RGBColor(0x24, 0x22, 0x1E)  # surface
ACCENT     = RGBColor(0xE8, 0x5D, 0x26)  # warm orange (Hive)
ACCENT2    = RGBColor(0x4E, 0xC5, 0xD1)  # teal accent
GREEN      = RGBColor(0x4A, 0xDE, 0x80)  # success green
WHITE      = RGBColor(0xF8, 0xF5, 0xF0)  # warm white
CREAM      = RGBColor(0xE8, 0xE0, 0xD4)  # warm gray
MUTED      = RGBColor(0x9C, 0x94, 0x88)  # muted text
CHART_COLORS = [ACCENT, ACCENT2, GREEN, RGBColor(0xF5, 0x9E, 0x0B), RGBColor(0xA7, 0x8B, 0xFA)]

FONT_HEADING = "Plus Jakarta Sans"
FONT_BODY    = "Inter"
FONT_MONO    = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]

# === HELPERS ===
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, w, h, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, w, h, text, font_size=14, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, font_name=FONT_BODY, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.2:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multiline(slide, left, top, w, h, lines, default_size=14, default_color=WHITE,
                  font_name=FONT_BODY, line_spacing=1.5):
    """lines = [(text, size, color, bold), ...] or just text strings."""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, sz, clr, bld = line, default_size, default_color, False
        else:
            txt = line[0]
            sz = line[1] if len(line) > 1 else default_size
            clr = line[2] if len(line) > 2 else default_color
            bld = line[3] if len(line) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(sz * 0.3)
        p.space_before = Pt(0)
    return txBox

def add_table(slide, left, top, w, h, rows, cols, data, header_bg=ACCENT,
              header_fg=BG_DARK, cell_bg=BG_CARD, cell_fg=WHITE, font_size=11):
    """data = [[cell_text, ...], ...] including header row."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, w, h)
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] is not None else ""
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = FONT_BODY
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = header_fg
                else:
                    p.font.color.rgb = cell_fg
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_bg if r % 2 == 1 else BG_SURFACE
            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
    return table_shape

def add_chart_bar(slide, left, top, w, h, categories, series_data, chart_title=""):
    """series_data = [(series_name, [values]), ...]"""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, vals in series_data:
        chart_data.add_series(name, vals)
    chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, w, h, chart_data)
    chart = chart_frame.chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.color.rgb = MUTED
    style_chart(chart)
    return chart_frame

def add_chart_line(slide, left, top, w, h, categories, series_data, chart_title=""):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, vals in series_data:
        chart_data.add_series(name, vals)
    chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, w, h, chart_data)
    chart = chart_frame.chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.color.rgb = MUTED
    style_chart(chart)
    return chart_frame

def add_chart_pie(slide, left, top, w, h, categories, values):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("", values)
    chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, w, h, chart_data)
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(9)
    chart.legend.font.color.rgb = MUTED
    style_chart(chart)
    return chart_frame

def style_chart(chart):
    """Apply dark theme to chart."""
    try:
        plot = chart.plots[0]
        plot.gap_width = 80
        # Color each series
        for i, series in enumerate(plot.series):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = CHART_COLORS[i % len(CHART_COLORS)]
    except:
        pass
    # Dark background
    chart.chart_style = 2
    try:
        chart.element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/chart}plotArea/{http://schemas.openxmlformats.org/drawingml/2006/chart}spPr')
    except:
        pass

def add_tag(slide, left, top, text, color=ACCENT):
    """Add a small tag/badge."""
    shape = add_rect(slide, left, top, Inches(1.8), Inches(0.32), RGBColor(
        min(255, color.red + 20) if hasattr(color, 'red') else color.red,
        min(255, color.green + 20) if hasattr(color, 'green') else color.green,
        min(255, color.blue + 20) if hasattr(color, 'blue') else color.blue,
    ), color, Pt(1))
    shape.fill.solid()
    # Semi-transparent effect via lighter shade
    shape.fill.fore_color.rgb = RGBColor(
        min(255, BG_DARK.red + 15),
        min(255, BG_DARK.green + 15),
        min(255, BG_DARK.blue + 15),
    )
    add_text(slide, left + Inches(0.1), top + Inches(0.02), Inches(1.6), Inches(0.28),
             text, font_size=9, color=color, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)

def add_metric_card(slide, left, top, w, h, label, value, delta=None, delta_color=GREEN):
    card = add_rect(slide, left, top, w, h, BG_CARD, RGBColor(0x33, 0x30, 0x2A), Pt(1))
    add_text(slide, left + Inches(0.2), top + Inches(0.15), w - Inches(0.4), Inches(0.25),
             label, font_size=9, color=MUTED, bold=True, font_name=FONT_MONO)
    add_text(slide, left + Inches(0.2), top + Inches(0.4), w - Inches(0.4), Inches(0.45),
             value, font_size=22, color=WHITE, bold=True, font_name=FONT_HEADING)
    if delta:
        add_text(slide, left + Inches(0.2), top + h - Inches(0.35), w - Inches(0.4), Inches(0.25),
                 delta, font_size=9, color=delta_color, bold=True, font_name=FONT_MONO)

# =========================================================
# SLIDE 1 — TITLE
# =========================================================
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)
# Subtle accent line
add_rect(s1, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
# Tag
add_tag(s1, Inches(1), Inches(1.8), "SERIES A  |  2026", ACCENT)
# Company name
add_text(s1, Inches(1), Inches(2.4), Inches(11), Inches(1.0),
         "BatteryHive", font_size=64, color=WHITE, bold=True, font_name=FONT_HEADING)
# Tagline
add_text(s1, Inches(1), Inches(3.5), Inches(9), Inches(0.6),
         "The future of EV energy is a swap away.", font_size=24, color=CREAM, font_name=FONT_HEADING)
# Sub-description
add_text(s1, Inches(1), Inches(4.3), Inches(8), Inches(0.8),
         "Automated battery swapping stations that replace depleted EV batteries in under 90 seconds.\nZero wait. Zero range anxiety. 100% renewable.",
         font_size=14, color=MUTED, line_spacing=1.6)
# Bottom metrics row
metrics = [
    ("STATIONS LIVE", "127"),
    ("SWAPS / DAY", "8,400+"),
    ("AVG SWAP TIME", "82 sec"),
    ("SERIES A ASK", "$45M"),
]
for i, (label, val) in enumerate(metrics):
    x = Inches(1 + i * 2.8)
    add_rect(s1, x, Inches(5.8), Inches(2.4), Inches(1.0), BG_CARD, RGBColor(0x33, 0x30, 0x2A), Pt(1))
    add_text(s1, x + Inches(0.15), Inches(5.9), Inches(2.1), Inches(0.2),
             label, font_size=8, color=MUTED, bold=True, font_name=FONT_MONO)
    add_text(s1, x + Inches(0.15), Inches(6.15), Inches(2.1), Inches(0.5),
             val, font_size=26, color=ACCENT, bold=True, font_name=FONT_HEADING)

# Confidential footer
add_text(s1, Inches(1), Inches(7.0), Inches(11), Inches(0.3),
         "CONFIDENTIAL  \u2022  BatteryHive Inc.  \u2022  July 2026",
         font_size=8, color=RGBColor(0x55, 0x50, 0x48), alignment=PP_ALIGN.LEFT, font_name=FONT_MONO)


# =========================================================
# SLIDE 2 — THE PROBLEM
# =========================================================
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)
add_rect(s2, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
add_tag(s2, Inches(0.8), Inches(0.4), "THE PROBLEM", ACCENT)
add_text(s2, Inches(0.8), Inches(0.9), Inches(6), Inches(0.6),
         "EV adoption is accelerating.\nCharging infrastructure isn\u2019t.", font_size=28, color=WHITE, bold=True, font_name=FONT_HEADING)

# Problem stats cards
problems = [
    ("38 MIN", "Average DC fast\ncharge to 80%", "\u2191 4.2x slower than gas"),
    ("31%", "EV owners cite charging\nas #1 pain point", "McKinsey 2025 Survey"),
    ("$18B", "Lost productivity from\ncharging wait time/yr", "US market estimate"),
    ("72%", "Urban EV owners lack\nhome charging access", "DOE 2025 Report"),
]
for i, (val, desc, sub) in enumerate(problems):
    x = Inches(0.8 + i * 3.05)
    card = add_rect(s2, x, Inches(2.0), Inches(2.8), Inches(2.0), BG_CARD, RGBColor(0x33, 0x30, 0x2A), Pt(1))
    add_text(s2, x + Inches(0.2), Inches(2.15), Inches(2.4), Inches(0.5),
             val, font_size=32, color=ACCENT, bold=True, font_name=FONT_HEADING)
    add_text(s2, x + Inches(0.2), Inches(2.75), Inches(2.4), Inches(0.6),
             desc, font_size=12, color=WHITE, line_spacing=1.4)
    add_text(s2, x + Inches(0.2), Inches(3.5), Inches(2.4), Inches(0.3),
             sub, font_size=9, color=MUTED, font_name=FONT_MONO)

# Charging comparison table
table_data = [
    ["Method", "Time to 80%", "Cost / 100mi", "Convenience"],
    ["Level 2 Home", "4\u20138 hours", "$4.50", "\u2605\u2605\u2605\u2605\u2605"],
    ["DC Fast Charge", "20\u201345 min", "$8\u201314", "\u2605\u2605\u2606"],
    ["Gas Station", "3\u20135 min", "$12\u201318", "\u2605\u2605\u2605\u2605"],
    ["BatteryHive Swap", "82 sec", "$9.50", "\u2605\u2605\u2605\u2605\u2605"],
]
add_table(s2, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.2),
          5, 4, table_data, font_size=11)

# Bottom insight
add_rect(s2, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.45), BG_SURFACE)
add_text(s2, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.35),
         "\u26a1  The gap: EVs need to be as fast as gas. BatteryHive makes them faster.",
         font_size=11, color=ACCENT, bold=True)


# =========================================================
# SLIDE 3 — THE SOLUTION
# =========================================================
s3 = prs.slides.add_slide(BLANK)
add_bg(s3)
add_rect(s3, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT2)
add_tag(s3, Inches(0.8), Inches(0.4), "THE SOLUTION", ACCENT2)
add_text(s3, Inches(0.8), Inches(0.9), Inches(6), Inches(0.6),
         "Swap, don\u2019t charge.\nDrive in, drive out.", font_size=28, color=WHITE, bold=True, font_name=FONT_HEADING)

# How it works steps
steps = [
    ("01", "ARRIVE", "Pull into any BatteryHive station. Camera-guided alignment \u2014 no precision needed."),
    ("02", "DIAGNOSE", "AI scans your battery health, charge level, and vehicle profile in 8 seconds."),
    ("03", "SWAP", "Robotic arm removes depleted pack, installs charged unit. 82-second average."),
    ("04", "GO", "Payment auto-processes. You drive away with 95% charge. No app needed."),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8 + i * 3.05)
    card = add_rect(s3, x, Inches(2.0), Inches(2.8), Inches(2.6), BG_CARD, RGBColor(0x33, 0x30, 0x2A), Pt(1))
    add_text(s3, x + Inches(0.2), Inches(2.15), Inches(0.6), Inches(0.4),
             num, font_size=28, color=ACCENT2, bold=True, font_name=FONT_MONO)
    add_text(s3, x + Inches(0.2), Inches(2.6), Inches(2.4), Inches(0.3),
             title, font_size=13, color=WHITE, bold=True, font_name=FONT_HEADING)
    add_text(s3, x + Inches(0.2), Inches(3.0), Inches(2.4), Inches(1.2),
             desc, font_size=11, color=CREAM, line_spacing=1.5)

# Key advantages
add_text(s3, Inches(0.8), Inches(5.0), Inches(4), Inches(0.3),
         "KEY ADVANTAGES", font_size=9, color=MUTED, bold=True, font_name=FONT_MONO)
advantages = [
    ("\u2713  10x faster than DC fast charging", GREEN),
    ("\u2713  Compatible with 85% of EV models via universal adapter", GREEN),
    ("\u2713  Battery health AI extends pack life by 23%", GREEN),
    ("\u2713  Station footprint: 2 parking spaces vs. 8 for charge lot", GREEN),
    ("\u2713  Grid-friendly: batteries charge off-peak, swap on-demand", GREEN),
]
for i, (txt, clr) in enumerate(advantages):
    add_text(s3, Inches(0.8), Inches(5.35 + i * 0.35), Inches(5), Inches(0.3),
             txt, font_size=11, color=clr)

# Tech stack card
card = add_rect(s3, Inches(7.5), Inches(5.0), Inches(5.0), Inches(2.2), BG_CARD, RGBColor(0x33, 0x30, 0x2A), Pt(1))
add_text(s3, Inches(7.7), Inches(5.1), Inches(4.6), Inches(0.25),
         "TECHNOLOGY STACK", font_size=9, color=MUTED, bold=True, font_name=FONT_MONO)
tech_lines = [
    ("Robotic Arm", "6-axis servo, <0.5mm precision, IP67 rated"),
    ("Battery Pack", "LFP 75 kWh modular, 2,500+ cycle life"),
    ("AI Diagnostics", "Real-time SOH/SOC via 48-cell BMS telemetry"),
    ("Grid Integration", "V2G bidirectional, 150kW charge / 350kW discharge"),
    ("Payments", "NFC + app + fleet API, zero-touch billing"),
]
for i, (label, val) in enumerate(tech_lines):
    add_text(s3, Inches(7.7), Inches(5.45 + i * 0.33), Inches(1.4), Inches(0.3),
             label, font_size=9, color=ACCENT2, bold=True)
    add_text(s3, Inches(9.2), Inches(5.45 + i * 0.33), Inches(3.1), Inches(0.3),
             val, font_size=9, color=CREAM)


# =========================================================
# SLIDE 4 — MARKET OPPORTUNITY
# =========================================================
s4 = prs.slides.add_slide(BLANK)
add_bg(s4)
add_rect(s4, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
add_tag(s4, Inches(0.8), Inches(0.4), "MARKET OPPORTUNITY", ACCENT)
add_text(s4, Inches(0.8), Inches(0.9), Inches(6), Inches(0.6),
         "A $340B market hiding\nin plain sight.", font_size=28, color=WHITE, bold=True, font_name=FONT_HEADING)

# TAM/SAM/SOM cards
tam_data = [
    ("TAM", "$340B", "Global EV charging &\nenergy infrastructure by 2030", ACCENT),
    ("SAM", "$52B", "Urban battery swap &\nfast-charge market, 12 countries", ACCENT2),
    ("SOM", "$3.8B", "Addressable by 2028 with\ncurrent expansion plan", GREEN),
]
for i, (label, val, desc, clr) in enumerate(tam_data):
    x = Inches(0.8 + i * 4.0)
    card = add_rect(s4, x, Inches(2.0), Inches(3.6), Inches(1.8), BG_CARD, RGBColor(0x33, 0x30, 0x2A), Pt(1))
    add_text(s4, x + Inches(0.2), Inches(2.1), Inches(3.2), Inches(0.25),
             label, font_size=10, color=clr, bold=True, font_name=FONT_MONO)
    add_text(s4, x + Inches(0.2), Inches(2.4), Inches(3.2), Inches(0.5),
             val, font_size=36, color=WHITE, bold=True, font_name=FONT_HEADING)
    add_text(s4, x + Inches(0.2), Inches(3.0), Inches(3.2), Inches(0.6),
             desc, font_size=11, color=CREAM, line_spacing=1.4)

# EV adoption chart
chart1 = add_chart_bar(s4, Inches(0.8), Inches(4.2), Inches(5.5), Inches(3.0),
    ["2022", "2023", "2024", "2025", "2026E", "2027E", "2028E"],
    [("Global EV Sales (M)", [10.5, 14.2, 17.8, 22.1, 28.5, 35.0, 42.0]),
     ("Battery Swap Ready (M)", [0.3, 0.8, 1.9, 4.2, 8.5, 15.0, 24.0])])

# Market growth chart
chart2 = add_chart_line(s4, Inches(7.0), Inches(4.2), Inches(5.5), Inches(3.0),
    ["2024", "2025", "2026", "2027", "2028", "2029", "2030"],
    [("Swap Station Revenue ($B)", [1.2, 2.8, 5.4, 9.8, 16.2, 24.5, 35.0]),
     ("Charging Infra Revenue ($B)", [18, 24, 31, 39, 48, 58, 68])])


# =========================================================
# SLIDE 5 — PRODUCT DEEP DIVE
# =========================================================
s5 = prs.slides.add_slide(BLANK)
add_bg(s5)
add_rect(s5, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT2)
add_tag(s5, Inches(0.8), Inches(0.4), "PRODUCT", ACCENT2)
add_text(s5, Inches(0.8), Inches(0.9), Inches(6), Inches(0.6),
         "The BatteryHive Station", font_size=28, color=WHITE, bold=True, font_name=FONT_HEADING)

# Station specs table
station_data = [
    ["Specification", "Value", "Industry Avg"],
    ["Swap Time (avg)", "82 seconds", "3\u20135 min (NIO)"],
    ["Battery Capacity", "75 kWh (LFP)", "70 kWh (NIO)"],
    ["Cycle Life", "2,500+ cycles", "1,500 cycles"],
    ["Charge Rate (off-peak)", "150 kW DC", "60\u2013150 kW"],
    ["Station Footprint", "2 parking spots", "8\u201312 spots"],
    ["Daily Throughput", "120 swaps/station", "40 charges/station"],
    ["Uptime", "99.4%", "92\u201396%"],
    ["Compatible Models", "85% of BEVs", "10\u201315%"],
]
add_table(s5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
          9, 3, station_data, font_size=10)

# Station cost breakdown
add_text(s5, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.3),
         "STATION UNIT ECONOMICS", font_size=9, color=MUTED, bold=True, font_name=FONT_MONO)

cost_data = [
    ["Component", "Cost", "% of Total"],
    ["Robotic Arm + Mechanism", "$185,000", "37%"],
    ["Battery Inventory (20 packs)", "$160,000", "32%"],
    ["Grid Connection + V2G", "$65,000", "13%"],
    ["Construction + Install", "$48,000", "10%"],
    ["AI / Software / Sensors", "$42,000", "8%"],
    ["Total Per Station", "$500,000", "100%"],
]
add_table(s5, Inches(7.0), Inches(2.2), Inches(5.5), Inches(3.2),
          7, 3, cost_data, font_size=10)

# ROI card
card = add_rect(s5, Inches(7.0), Inches(5.6), Inches(5.5), Inches(1.5), BG_CARD, RGBColor(0x33, 0x30, 0x2A), Pt(1))
add_text(s5, Inches(7.2), Inches(5.7), Inches(5.1), Inches(0.25),
         "STATION ROI", font_size=9, color=GREEN, bold=True, font_name=FONT_MONO)
roi_lines = [
    ("Revenue / day (80 swaps \u00d7 $9.50)", "$760", WHITE),
    ("Monthly revenue", "$22,800", WHITE),
    ("Monthly opex (energy + maintenance)", "$8,400", MUTED),
    ("Monthly margin", "$14,400", GREEN),
    ("Payback period", "34 months", ACCENT),
]
for i, (label, val, clr) in enumerate(roi_lines):
    add_text(s5, Inches(7.2), Inches(6.0 + i * 0.22), Inches(3.5), Inches(0.2),
             label, font_size=9, color=MUTED)
    add_text(s5, Inches(10.8), Inches(6.0 + i * 0.22), Inches(1.5), Inches(0.2),
             val, font_size=9, color=clr, bold=True, alignment=PP_ALIGN.RIGHT)


# =========================================================
# SLIDE 6 — BUSINESS MODEL & REVENUE
# =========================================================
s6 = prs.slides.add_slide(BLANK)
add_bg(s6)
add_rect(s6, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
add_tag(s6, Inches(0.8), Inches(0.4), "BUSINESS MODEL", ACCENT)
add_text(s6, Inches(0.8), Inches(0.9), Inches(6), Inches(0.6),
         "Three revenue streams.\nOne platform.", font_size=28, color=WHITE, bold=True, font_name=FONT_HEADING)

# Revenue streams
streams = [
    ("PER-SWAP FEE", "$9.50 / swap", "62% of revenue",
     "Pay-per-use pricing. Average user swaps 3.2x/week.\nFleet contracts at $8.20/swap (volume discount).", ACCENT),
    ("FLEET CONTRACTS", "$2,400 / mo", "28% of revenue",
     "B2B subscriptions for ride-share, delivery, logistics.\n12-month minimum. Includes priority lane access.", ACCENT2),
    ("GRID SERVICES", "$1,800 / mo", "10% of revenue",
     "V2G peak shaving, demand response, frequency regulation.\nBatteries as grid assets during off-swap hours.", GREEN),
]
for i, (title, val, pct, desc, clr) in enumerate(streams):
    x = Inches(0.8 + i * 4.0)
    card = add_rect(s6, x, Inches(2.0), Inches(3.6), Inches(2.8), BG_CARD, RGBColor(0x33, 0x30, 0x2A), Pt(1))
    add_text(s6, x + Inches(0.2), Inches(2.1), Inches(3.2), Inches(0.25),
             title, font_size=9, color=clr, bold=True, font_name=FONT_MONO)
    add_text(s6, x + Inches(0.2), Inches(2.4), Inches(3.2), Inches(0.45),
             val, font_size=26, color=WHITE, bold=True, font_name=FONT_HEADING)
    add_text(s6, x + Inches(0.2), Inches(2.9), Inches(3.2), Inches(0.25),
             pct, font_size=10, color=clr)
    add_text(s6, x + Inches(0.2), Inches(3.3), Inches(3.2), Inches(1.0),
             desc, font_size=10, color=CREAM, line_spacing=1.5)

# Revenue breakdown chart
chart3 = add_chart_pie(s6, Inches(0.8), Inches(5.0), Inches(4.5), Inches(2.3),
    ["Per-Swap Fees", "Fleet Contracts", "Grid Services"],
    [62, 28, 10])

# Revenue projection chart
chart4 = add_chart_bar(s6, Inches(6.0), Inches(5.0), Inches(6.5), Inches(2.3),
    ["2026", "2027", "2028", "2029", "2030"],
    [("Swap Revenue ($M)", [8.2, 28.5, 68.0, 142.0, 285.0]),
     ("Fleet Revenue ($M)", [3.8, 12.0, 32.0, 68.0, 135.0]),
     ("Grid Revenue ($M)", [1.2, 4.5, 12.0, 28.0, 52.0])])


# =========================================================
# SLIDE 7 — TRACTION & METRICS
# =========================================================
s7 = prs.slides.add_slide(BLANK)
add_bg(s7)
add_rect(s7, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GREEN)
add_tag(s7, Inches(0.8), Inches(0.4), "TRACTION", GREEN)
add_text(s7, Inches(0.8), Inches(0.9), Inches(6), Inches(0.6),
         "We\u2019re not pitching a concept.\nWe\u2019re scaling a system.", font_size=28, color=WHITE, bold=True, font_name=FONT_HEADING)

# KPI cards row 1
kpis = [
    ("STATIONS LIVE", "127", "+34 this quarter", GREEN),
    ("TOTAL SWAPS", "1.82M", "since Jan 2025 launch", GREEN),
    ("MONTHLY SWAPS", "252K", "+18% MoM growth", ACCENT),
    ("ACTIVE USERS", "34,200", "92% monthly retention", ACCENT2),
]
for i, (label, val, sub, clr) in enumerate(kpis):
    x = Inches(0.8 + i * 3.05)
    add_metric_card(s7, x, Inches(1.8), Inches(2.8), Inches(1.5), label, val, sub, clr)

# Monthly swaps chart
chart5 = add_chart_line(s7, Inches(0.8), Inches(3.6), Inches(5.5), Inches(3.5),
    ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"],
    [("2025 Monthly Swaps (K)", [42, 51, 58, 68, 79, 92, 108, 118, 132, 148, 168, 195]),
     ("2026 Monthly Swaps (K)", [198, 215, 228, 242, 252, 268, 285, 305, 328, 355, 385, 420])])

# Unit economics table
ue_data = [
    ["Metric", "Current", "Target (2027)", "Industry"],
    ["Customer Acq. Cost", "$18", "$12", "$35"],
    ["LTV (24-month)", "$1,420", "$1,850", "$680"],
    ["LTV:CAC Ratio", "79x", "154x", "19x"],
    ["Monthly Churn", "1.2%", "0.8%", "3.5%"],
    ["Net Revenue Retention", "118%", "128%", "95%"],
    ["Gross Margin", "68%", "74%", "42%"],
]
add_table(s7, Inches(7.0), Inches(3.6), Inches(5.5), Inches(3.3),
          7, 4, ue_data, font_size=10)


# =========================================================
# SLIDE 8 — TEAM
# =========================================================
s8 = prs.slides.add_slide(BLANK)
add_bg(s8)
add_rect(s8, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT2)
add_tag(s8, Inches(0.8), Inches(0.4), "THE TEAM", ACCENT2)
add_text(s8, Inches(0.8), Inches(0.9), Inches(6), Inches(0.6),
         "Operators who\u2019ve built\nat scale before.", font_size=28, color=WHITE, bold=True, font_name=FONT_HEADING)

team = [
    ("Sarah Chen", "CEO & Co-Founder",
     "Ex-Tesla VP, Supercharger Network. Built 12,000+ stations globally. Stanford MS EE. 3 patents in fast-charging.",
     "Tesla \u2022 Stanford \u2022 3 Patents"),
    ("Marcus Williams", "CTO & Co-Founder",
     "Ex-Boston Dynamics, lead robotics engineer. Built autonomous charging systems. MIT PhD Mechanical Engineering.",
     "Boston Dynamics \u2022 MIT PhD"),
    ("Dr. Amara Osei", "VP Battery Engineering",
     "Ex-Panasonic, designed cells for Tesla Gigafactory. 15 years in lithium-ion. 8 patents.",
     "Panasonic \u2022 8 Patents"),
    ("James Park", "VP Operations",
     "Ex-Lyft, scaled ops from 5 to 200 cities. Harvard MBA. Expert in marketplace logistics.",
     "Lyft \u2022 Harvard MBA"),
]
for i, (name, role, bio, creds) in enumerate(team):
    x = Inches(0.8 + i * 3.05)
    card = add_rect(s8, x, Inches(2.0), Inches(2.8), Inches(3.5), BG_CARD, RGBColor(0x33, 0x30, 0x2A), Pt(1))
    # Avatar placeholder
    avatar = add_rect(s8, x + Inches(0.2), Inches(2.15), Inches(0.7), Inches(0.7), BG_SURFACE, ACCENT2, Pt(1))
    add_text(s8, x + Inches(0.2), Inches(2.25), Inches(0.7), Inches(0.5),
             name[0] + name.split()[-1][0], font_size=16, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)
    add_text(s8, x + Inches(1.0), Inches(2.2), Inches(1.6), Inches(0.3),
             name, font_size=13, color=WHITE, bold=True, font_name=FONT_HEADING)
    add_text(s8, x + Inches(1.0), Inches(2.5), Inches(1.6), Inches(0.25),
             role, font_size=9, color=ACCENT2, font_name=FONT_MONO)
    add_text(s8, x + Inches(0.2), Inches(3.1), Inches(2.4), Inches(1.5),
             bio, font_size=10, color=CREAM, line_spacing=1.5)
    add_text(s8, x + Inches(0.2), Inches(4.8), Inches(2.4), Inches(0.25),
             creds, font_size=8, color=MUTED, font_name=FONT_MONO)

# Advisors
add_text(s8, Inches(0.8), Inches(5.8), Inches(2), Inches(0.3),
         "ADVISORS", font_size=9, color=MUTED, bold=True, font_name=FONT_MONO)
advisors = [
    "Dr. JB Straubel  \u2014  Co-Founder, Tesla & Redwood Materials",
    "Andrew Chen  \u2014  General Partner, a16z",
    "Megan Smith  \u2014  Former US CTO, Google VP",
]
for i, adv in enumerate(advisors):
    add_text(s8, Inches(0.8), Inches(6.15 + i * 0.32), Inches(11), Inches(0.3),
             adv, font_size=10, color=CREAM)


# =========================================================
# SLIDE 9 — FINANCIAL PROJECTIONS
# =========================================================
s9 = prs.slides.add_slide(BLANK)
add_bg(s9)
add_rect(s9, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
add_tag(s9, Inches(0.8), Inches(0.4), "FINANCIALS", ACCENT)
add_text(s9, Inches(0.8), Inches(0.9), Inches(6), Inches(0.6),
         "Path to $1B revenue\nby 2030.", font_size=28, color=WHITE, bold=True, font_name=FONT_HEADING)

# Financial table
fin_data = [
    ["", "2026E", "2027E", "2028E", "2029E", "2030E"],
    ["Stations Deployed", "127", "380", "850", "1,600", "3,200"],
    ["Total Swaps (M)", "1.8", "6.2", "18.5", "48.0", "112.0"],
    ["Revenue ($M)", "$13.2", "$45.0", "$112.0", "$238.0", "$472.0"],
    ["COGS ($M)", "$4.2", "$12.6", "$29.0", "$57.0", "$104.0"],
    ["Gross Margin", "68%", "72%", "74%", "76%", "78%"],
    ["OpEx ($M)", "$22.0", "$38.0", "$62.0", "$95.0", "$138.0"],
    ["EBITDA ($M)", "($13.0)", "($5.6)", "$21.0", "$86.0", "$230.0"],
    ["EBITDA Margin", "n/a", "n/a", "19%", "36%", "49%"],
    ["Cash Burn ($M/mo)", "$1.8", "$1.2", "$0.4", "Cash +", "Cash ++"],
]
add_table(s9, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8),
          10, 6, fin_data, font_size=10)

# Revenue + EBITDA chart
chart6 = add_chart_bar(s9, Inches(8.8), Inches(1.8), Inches(4.0), Inches(2.8),
    ["2026E", "2027E", "2028E", "2029E", "2030E"],
    [("Revenue ($M)", [13.2, 45.0, 112.0, 238.0, 472.0]),
     ("EBITDA ($M)", [-13.0, -5.6, 21.0, 86.0, 230.0])])

# Station growth chart
chart7 = add_chart_line(s9, Inches(8.8), Inches(4.8), Inches(4.0), Inches(2.3),
    ["2026E", "2027E", "2028E", "2029E", "2030E"],
    [("Cumulative Stations", [127, 380, 850, 1600, 3200])])

# Key assumptions
add_text(s9, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
         "Key assumptions: 80 swaps/station/day avg, $9.50/swap blended, 78% gross margin at scale, 18\u201324 month payback per station.",
         font_size=9, color=MUTED, font_name=FONT_MONO)


# =========================================================
# SLIDE 10 — THE ASK
# =========================================================
s10 = prs.slides.add_slide(BLANK)
add_bg(s10)
add_rect(s10, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
add_tag(s10, Inches(0.8), Inches(0.4), "THE ASK", ACCENT)
add_text(s10, Inches(0.8), Inches(0.9), Inches(6), Inches(0.6),
         "$45M Series A\nto dominate the swap economy.", font_size=28, color=WHITE, bold=True, font_name=FONT_HEADING)

# Use of funds
add_text(s10, Inches(0.8), Inches(2.0), Inches(4), Inches(0.3),
         "USE OF FUNDS", font_size=9, color=MUTED, bold=True, font_name=FONT_MONO)

funds = [
    ("Station Deployment", "$22M", "49%", "200 new stations in 8 metro areas"),
    ("R&D / Battery Tech", "$10M", "22%", "Next-gen solid-state pack development"),
    ("Team Expansion", "$7M", "16%", "80 new hires across eng, ops, sales"),
    ("Working Capital", "$4M", "9%", "Battery inventory for new stations"),
    ("Regulatory / Legal", "$2M", "4%", "Permits, certifications, IP protection"),
]
for i, (item, amount, pct, desc) in enumerate(funds):
    y = Inches(2.4 + i * 0.82)
    card = add_rect(s10, Inches(0.8), y, Inches(5.5), Inches(0.7), BG_CARD, RGBColor(0x33, 0x30, 0x2A), Pt(1))
    add_text(s10, Inches(1.0), y + Inches(0.08), Inches(2.5), Inches(0.25),
             item, font_size=11, color=WHITE, bold=True)
    add_text(s10, Inches(3.5), y + Inches(0.08), Inches(1.0), Inches(0.25),
             amount, font_size=11, color=ACCENT, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text(s10, Inches(4.6), y + Inches(0.08), Inches(0.5), Inches(0.25),
             pct, font_size=9, color=MUTED, alignment=PP_ALIGN.RIGHT)
    add_text(s10, Inches(1.0), y + Inches(0.38), Inches(5.1), Inches(0.25),
             desc, font_size=9, color=CREAM)
    # Progress bar
    bar_bg = add_rect(s10, Inches(5.5), y + Inches(0.12), Inches(0.6), Inches(0.15), BG_SURFACE)
    bar_fg = add_rect(s10, Inches(5.5), y + Inches(0.12),
                      Inches(0.6 * float(pct.replace('%','')) / 100), Inches(0.15), ACCENT)

# Use of funds pie
chart8 = add_chart_pie(s10, Inches(7.0), Inches(2.0), Inches(5.5), Inches(2.8),
    ["Station Deployment", "R&D / Battery", "Team", "Working Capital", "Regulatory"],
    [49, 22, 16, 9, 4])

# Milestones
add_text(s10, Inches(7.0), Inches(5.0), Inches(5), Inches(0.3),
         "18-MONTH MILESTONES", font_size=9, color=GREEN, bold=True, font_name=FONT_MONO)
milestones = [
    ("Q4 2026", "250 stations live, 50K daily swaps"),
    ("Q2 2027", "Break-even on unit economics, 10 metro areas"),
    ("Q4 2027", "500 stations, $45M ARR, Series B ready"),
]
for i, (q, desc) in enumerate(milestones):
    y = Inches(5.4 + i * 0.45)
    add_rect(s10, Inches(7.0), y, Inches(0.08), Inches(0.3), GREEN)
    add_text(s10, Inches(7.3), y + Inches(0.02), Inches(1.2), Inches(0.25),
             q, font_size=11, color=GREEN, bold=True, font_name=FONT_MONO)
    add_text(s10, Inches(8.6), y + Inches(0.02), Inches(3.7), Inches(0.25),
             desc, font_size=11, color=CREAM)

# Closing CTA
add_rect(s10, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), BG_SURFACE, ACCENT, Pt(1))
add_text(s10, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.4),
         "BatteryHive \u2014 The future of EV energy is a swap away.  \u2022  contact@batteryhive.com  \u2022  batteryhive.com",
         font_size=11, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

# === SAVE ===
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "BatteryHive_Pitch_Deck.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
